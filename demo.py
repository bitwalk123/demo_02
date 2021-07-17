#!/usr/bin/env python
# coding: utf-8
from PySide6.QtWidgets import (
    QApplication,
    QComboBox,
    QFileDialog,
    QMainWindow,
    QMessageBox,
    QSizePolicy,
    QStatusBar,
    QToolBar,
    QToolButton,
    QWidget,
)
import pandas as pd
import sys
import tempfile

import pptx
from pptx import Presentation, shapes, slide
from pptx.util import Inches

# local python files
import demo_common as cm
from demo_chart import TrendChart
from demo_database import DBObj
from demo_dialog import AboutDlg
from demo_icon import AppIcon
from demo_parser import FileParser


class Demo(QMainWindow):
    # Application Information
    APP_NAME: str = 'Demo 2'
    APP_VER: str = '0.1'
    APP_COPYRIGHT: str = '2021 Fuhito Suguri'
    APP_LICENSE: str = '<a href="https://opensource.org/licenses/MIT">MIT</a>'

    # GUI related
    w_init: int = 800
    h_init: int = 600
    width_combo: int = 150

    # Database
    db_obj: DBObj = None
    name_db: str = 'sample.sqlite'

    # PowerPoint template
    template = 'template.pptx'

    def __init__(self):
        super().__init__()
        self.combo = QComboBox()
        self.init_ui()

        self.setWindowIcon(AppIcon('app'))
        self.setWindowTitle(self.APP_NAME)
        self.resize(self.w_init, self.h_init)

    # -------------------------------------------------------------------------
    #  init_ui
    #  initialize UI
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def init_ui(self):
        # Create Qt toolbar
        toolbar = QToolBar()
        self.addToolBar(toolbar)

        # open button
        but_open = QToolButton()
        but_open.setIcon(AppIcon('folder'))
        but_open.setStatusTip('Open CSV file')
        but_open.clicked.connect(self.open_file)
        toolbar.addWidget(but_open)
        toolbar.addSeparator()

        # combobox for parameters
        self.combo.setStatusTip('Select parameter to generate chart')
        self.combo.currentIndexChanged.connect(self.status_changed_combo)
        self.combo.setFixedWidth(self.width_combo)
        toolbar.addWidget(self.combo)
        toolbar.addSeparator()

        # excel button
        but_excel = QToolButton()
        but_excel.setIcon(AppIcon('excel'))
        but_excel.setStatusTip('Open Data with Excel file')
        but_excel.clicked.connect(self.excel_file)
        toolbar.addWidget(but_excel)

        # powerpoint button
        but_ppt = QToolButton()
        but_ppt.setIcon(AppIcon('powerpoint'))
        but_ppt.setStatusTip('Open Trend with PowerPoint file')
        but_ppt.clicked.connect(self.ppt_file)
        toolbar.addWidget(but_ppt)

        # spacer
        spacer: QWidget = QWidget()
        spacer.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        toolbar.addWidget(spacer)

        # info button
        but_info = QToolButton()
        but_info.setIcon(AppIcon('info'))
        but_info.setStatusTip('About this application')
        but_info.clicked.connect(self.about_dialog)
        toolbar.addWidget(but_info)

        # exit button
        but_exit = QToolButton()
        but_exit.setIcon(AppIcon('exit'))
        but_exit.setStatusTip('Exit this app')
        but_exit.clicked.connect(self.closeEvent)
        toolbar.addWidget(but_exit)

        # Create Qt statusbar
        statusbar = QStatusBar()
        self.setStatusBar(statusbar)

    # -------------------------------------------------------------------------
    #  closeEvent (override)
    #  Dialog for close confirmation
    #
    #  argument
    #    event
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def about_dialog(self, event):
        about = AboutDlg(self)
        about.show()

    # -------------------------------------------------------------------------
    #  closeEvent (override)
    #  Dialog for close confirmation
    #
    #  argument
    #    event
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def closeEvent(self, event):
        sender = self.sender()

        reply: QMessageBox.StandardButton = QMessageBox.warning(
            self,
            'Quit App',
            'Are you sure you want to quit?',
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No
        )

        if reply == QMessageBox.Yes:
            QApplication.exit(0)

    # -------------------------------------------------------------------------
    #  excel_file
    #  create excel file
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def excel_file(self):
        if self.db_obj is None:
            return

        # generate temporary filename for Excel
        save_path: str = tempfile.NamedTemporaryFile(suffix='.xlsx').name

        # get all data from database
        df: pd.DataFrame = self.db_obj.get_df_all()

        # generate Excel file
        df.to_excel(save_path, index=False)

        # open file with default application
        cm.open_file_with_default_app(save_path)

    # -------------------------------------------------------------------------
    #  open_file
    #  Open file dialog
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def open_file(self):
        # Filter for file extensions to read
        filters: str = 'CSV file (*.csv);; All (*.*)'

        # file selection dialog
        dialog = QFileDialog()
        dialog.setNameFilter(filters)

        if not dialog.exec_():
            return

        # read selected file
        filename: str = dialog.selectedFiles()[0]
        file_obj = FileParser(filename)
        self.db_obj = DBObj(self.name_db, file_obj)
        self.combo.addItems(self.db_obj.get_params())

    # -------------------------------------------------------------------------
    #  ppt_file
    #  create powerpoint file
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def ppt_file(self):
        if self.db_obj is None:
            return

        # PowerPoint instance
        ppt = Presentation(self.template)

        # generate temporary filename for Excel
        save_path: str = tempfile.NamedTemporaryFile(suffix='.pptx').name

        # list of parameters
        list_param: list = self.db_obj.get_params()
        for param in list_param:
            self.gen_ppt_slide(ppt, param)

        # save PowerPoint file
        ppt.save(save_path)

        # open file with default application
        cm.open_file_with_default_app(save_path)

    # -------------------------------------------------------------------------
    #  gen_ppt_slide
    #  create powerpoint slide
    #
    #  argument
    #    param: str - plot parameter name
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def gen_ppt_slide(self, ppt: Presentation, param: str):
        # get dataframe for specified parameter
        #df: pd.DataFrame = self.db_obj.get_trend_data(param)

        # generate trend chart
        #chart = TrendChart(df)
        #fig: matplotlib.figure.Figure = chart.gen_chart()

        # save chart as PNG image
        #image_path: str = tempfile.NamedTemporaryFile(suffix='.png').name
        #fig.savefig(image_path)

        # refer layout from master
        slide_layout: pptx.slide.SlideLayout = ppt.slide_layouts[0]
        slide: pptx.slide.Slide = ppt.slides.add_slide(slide_layout)
        # shapes: pptx.shapes.shapetree.SlideShapes = slide.shapes
        shapes: pptx.shapes.shapetree.SlideShapes = slide.shapes

        # slide title
        shapes.title.text = param

        # insert image
        ileft: Inches = Inches(0)
        itop: Inches = Inches(1.35)
        iheight: Inches = Inches(3.5)
        #slide.shapes.add_picture(image_path, left=ileft, top=itop, height=iheight)

    # -------------------------------------------------------------------------
    #  status_changed_combo
    #  combobox status changed
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def status_changed_combo(self):
        combo: QComboBox = self.sender()
        param: str = combo.currentText()
        df: pd.DataFrame = self.db_obj.get_trend_data(param)
        trend = TrendChart(df)
        self.setCentralWidget(trend)
        # canvas: matplotlib.backends.backend_qt5agg.FigureCanvasQTAgg = trend.get_canvas()
        # self.setCentralWidget(canvas)


def main():
    app = QApplication(sys.argv)
    demo = Demo()
    demo.show()
    sys.exit(app.exec())


if __name__ == '__main__':
    main()
